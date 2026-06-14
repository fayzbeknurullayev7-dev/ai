"""Slayd (.pptx) generatsiya yadrosi.

Ikki bosqich:
  1. `generate_outline(prompt)` — promptdan slaydlar rejasini (sarlavha +
     punktlar) quradi. Groq (llama-3.3) bilan JSON reja so'raydi; API key yo'q,
     internet yo'q yoki javob buzuq bo'lsa — OFFLINE heuristik rejaga qaytadi
     (ilova hech qachon yiqilmaydi, testlar keysiz o'tadi).
  2. `build_pptx(deck)` — python-pptx bilan rejani .pptx baytlariga aylantiradi.

Dizayn image_agent bilan bir xil: tashqi xizmat ixtiyoriy, asosiy oqim
har doim ishlaydi.
"""
from __future__ import annotations

import json
import re
from dataclasses import dataclass, field
from io import BytesIO
from typing import List, Optional

from app.core.config import settings

# Groq importi modul yuklanishini sekinlashtirmasligi uchun funksiya ichida.

_MODEL = "llama-3.3-70b-versatile"

_SYSTEM = (
    "Sen taqdimot (slayd) tuzuvchi yordamchisan. Foydalanuvchi mavzusiga ko'ra "
    "professional slaydlar rejasini tuzasan. FAQAT quyidagi JSON formatida javob "
    "ber, boshqa matn yozma:\n"
    '{"title": "Taqdimot sarlavhasi", "slides": [{"title": "Slayd sarlavhasi", '
    '"bullets": ["punkt 1", "punkt 2", "punkt 3"]}]}\n'
    "6-9 ta slayd, har birida 3-5 ta qisqa punkt. Sarlavha va punktlar o'zbek "
    "tilida (foydalanuvchi boshqa til so'ramasa)."
)

_MAX_SLIDES = 12
_MAX_BULLETS = 8


@dataclass
class Slide:
    title: str
    bullets: List[str] = field(default_factory=list)


@dataclass
class SlideDeck:
    title: str
    slides: List[Slide] = field(default_factory=list)

    @property
    def slide_count(self) -> int:
        # Titul slayd + kontent slaydlar.
        return 1 + len(self.slides)


# --------------------------------------------------------------------------- #
# Reja yaratish                                                               #
# --------------------------------------------------------------------------- #
async def generate_outline(prompt: str, *, client: object | None = None) -> SlideDeck:
    """Promptdan slaydlar rejasini quradi (Groq → fallback offline)."""
    prompt = (prompt or "").strip()
    if not prompt:
        prompt = "Yangi taqdimot"

    deck = await _try_groq_outline(prompt, client)
    if deck is not None and deck.slides:
        return deck
    return _offline_outline(prompt)


async def _try_groq_outline(prompt: str, client: object | None) -> Optional[SlideDeck]:
    """Groq bilan JSON reja so'raydi. Muvaffaqiyatsiz bo'lsa None qaytaradi."""
    if client is None:
        if not settings.GROQ_API_KEY:
            return None
        try:
            from groq import AsyncGroq

            client = AsyncGroq(api_key=settings.GROQ_API_KEY)
        except Exception:
            return None
    try:
        response = await client.chat.completions.create(  # type: ignore[attr-defined]
            model=_MODEL,
            messages=[
                {"role": "system", "content": _SYSTEM},
                {"role": "user", "content": prompt},
            ],
            max_tokens=2048,
            temperature=0.4,
        )
        content = response.choices[0].message.content or ""
        return _parse_outline(content)
    except Exception:
        return None


def _parse_outline(content: str) -> Optional[SlideDeck]:
    """LLM matnidan JSON rejani ajratib oladi (markdown fence'larga chidamli)."""
    raw = content.strip()
    # ```json ... ``` yoki ``` ... ``` bloklarini tozalaymiz.
    fence = re.search(r"```(?:json)?\s*(.*?)```", raw, re.DOTALL)
    if fence:
        raw = fence.group(1).strip()
    else:
        # Birinchi { dan oxirgi } gacha.
        start, end = raw.find("{"), raw.rfind("}")
        if start != -1 and end != -1 and end > start:
            raw = raw[start : end + 1]
    try:
        data = json.loads(raw)
    except (json.JSONDecodeError, ValueError):
        return None
    if not isinstance(data, dict):
        return None

    title = str(data.get("title") or "Taqdimot").strip()
    slides: List[Slide] = []
    for item in (data.get("slides") or [])[:_MAX_SLIDES]:
        if not isinstance(item, dict):
            continue
        s_title = str(item.get("title") or "").strip()
        bullets = [
            str(b).strip()
            for b in (item.get("bullets") or [])
            if str(b).strip()
        ][:_MAX_BULLETS]
        if s_title or bullets:
            slides.append(Slide(title=s_title or "Slayd", bullets=bullets))
    if not slides:
        return None
    return SlideDeck(title=title, slides=slides)


def _offline_outline(prompt: str) -> SlideDeck:
    """Internet/Groq bo'lmaganda promptdan deterministik oddiy reja."""
    topic = prompt.replace("\n", " ").strip()
    short = topic if len(topic) <= 60 else topic[:60].rstrip() + "…"
    title = short[:1].upper() + short[1:] if short else "Taqdimot"
    slides = [
        Slide(
            title="Kirish",
            bullets=[
                f"Mavzu: {short}",
                "Ushbu taqdimot mavzuning asosiy jihatlarini yoritadi.",
                "Maqsad va kutilayotgan natijalar.",
            ],
        ),
        Slide(
            title="Asosiy g'oyalar",
            bullets=[
                "Birinchi muhim jihat.",
                "Ikkinchi muhim jihat.",
                "Uchinchi muhim jihat.",
            ],
        ),
        Slide(
            title="Tafsilotlar",
            bullets=[
                "Amaliy misollar va tushuntirishlar.",
                "Afzalliklar va cheklovlar.",
                "E'tiborga olinadigan nuqtalar.",
            ],
        ),
        Slide(
            title="Xulosa",
            bullets=[
                "Asosiy fikrlar qisqacha.",
                "Keyingi qadamlar.",
                "Savol-javob.",
            ],
        ),
    ]
    return SlideDeck(title=title, slides=slides)


# --------------------------------------------------------------------------- #
# .pptx qurish                                                                #
# --------------------------------------------------------------------------- #
def build_pptx(deck: SlideDeck) -> bytes:
    """Rejani python-pptx bilan .pptx faylga (baytlar) aylantiradi."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    # 1) Titul slayd (layout 0: Title Slide).
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = deck.title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Nexus AI tomonidan yaratilgan"

    # 2) Kontent slaydlar (layout 1: Title and Content).
    content_layout = prs.slide_layouts[1]
    for slide in deck.slides:
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = slide.title
        body = s.placeholders[1].text_frame
        body.clear()
        bullets = slide.bullets or [slide.title]
        for i, bullet in enumerate(bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = bullet
            para.level = 0
            para.font.size = Pt(18)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
